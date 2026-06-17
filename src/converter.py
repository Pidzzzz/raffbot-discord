import subprocess
import shutil
from pathlib import Path
from PIL import Image
import img2pdf
import fitz

LIBRE_PATHS = [
    "C:\\Program Files\\LibreOffice\\program\\soffice.exe",
    "C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe",
    shutil.which("soffice"),
    shutil.which("libreoffice"),
]

LIBRE_PATH = None
for p in LIBRE_PATHS:
    if p and Path(p).exists():
        LIBRE_PATH = p
        break

def _libre_convert(input_path: str, output_dir: str) -> bool:
    if not LIBRE_PATH:
        return False
    try:
        result = subprocess.run(
            [LIBRE_PATH, "--headless", "--convert-to", "pdf",
             "--outdir", output_dir, input_path],
            capture_output=True,
            timeout=120,
        )
        return result.returncode == 0
    except Exception:
        return False

async def image_to_pdf(input_path: str, output_path: str) -> str | None:
    try:
        img = Image.open(input_path)
        if img.mode in ("RGBA", "P"):
            img = img.convert("RGB")
        with open(output_path, "wb") as f:
            f.write(img2pdf.convert(input_path))
        return output_path
    except Exception as e:
        print(f"[image_to_pdf] {e}")
        return None

async def images_to_pdf(input_paths: list[str], output_path: str) -> str | None:
    try:
        with open(output_path, "wb") as f:
            f.write(img2pdf.convert(input_paths))
        return output_path
    except Exception as e:
        print(f"[images_to_pdf] {e}")
        return None

async def docx_to_pdf(input_path: str, output_path: str) -> str | None:
    if LIBRE_PATH:
        out_dir = str(Path(output_path).parent)
        if _libre_convert(input_path, out_dir):
            expected = Path(out_dir) / f"{Path(input_path).stem}.pdf"
            if expected.exists():
                expected.rename(output_path)
                return output_path
    try:
        from docx import Document as DocxDocument
        doc = DocxDocument(input_path)
        pdf = fitz.open()
        page = pdf.new_page()
        y, margin, lh = 50, 50, 14

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                y += lh
                continue
            fs = 18 if para.style and "Heading 1" in para.style.name else (
                14 if para.style and "Heading" in para.style.name else 12)
            if y + lh > page.rect.height - margin:
                page = pdf.new_page(); y = 50
            page.insert_text((margin, y), text, fontsize=fs)
            y += int(lh * (1.5 if fs > 12 else 1))

        for table in doc.tables:
            for row in table.rows:
                if y + lh > page.rect.height - margin:
                    page = pdf.new_page(); y = 50
                cells = [cell.text.strip() for cell in row.cells]
                page.insert_text((margin, y), " | ".join(cells), fontsize=10)
                y += lh
            y += 10

        pdf.save(output_path)
        pdf.close()
        return output_path
    except Exception as e:
        print(f"[docx_to_pdf] {e}")
        return None

async def xlsx_to_pdf(input_path: str, output_path: str) -> str | None:
    if LIBRE_PATH:
        out_dir = str(Path(output_path).parent)
        if _libre_convert(input_path, out_dir):
            expected = Path(out_dir) / f"{Path(input_path).stem}.pdf"
            if expected.exists():
                expected.rename(output_path)
                return output_path
    try:
        from openpyxl import load_workbook
        wb = load_workbook(input_path, data_only=True)
        pdf = fitz.open()
        page = pdf.new_page()
        y, margin, lh = 50, 50, 14

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            if y + lh * 2 > page.rect.height - margin:
                page = pdf.new_page(); y = 50
            page.insert_text((margin, y), f"[ {sheet_name} ]", fontsize=14, color=(0, 0, 0.6))
            y += lh * 2

            for row in ws.iter_rows():
                if y + lh > page.rect.height - margin:
                    page = pdf.new_page(); y = 50
                vals = [str(c.value) if c.value is not None else "" for c in row]
                line = " | ".join(vals)
                if line.strip():
                    page.insert_text((margin, y), line, fontsize=10)
                    y += lh
            y += 20

        pdf.save(output_path)
        pdf.close()
        wb.close()
        return output_path
    except Exception as e:
        print(f"[xlsx_to_pdf] {e}")
        return None

async def pptx_to_pdf(input_path: str, output_path: str) -> str | None:
    if LIBRE_PATH:
        out_dir = str(Path(output_path).parent)
        if _libre_convert(input_path, out_dir):
            expected = Path(out_dir) / f"{Path(input_path).stem}.pdf"
            if expected.exists():
                expected.rename(output_path)
                return output_path
    try:
        from pptx import Presentation
        prs = Presentation(input_path)
        pdf = fitz.open()

        for slide_num, slide in enumerate(prs.slides, 1):
            page = pdf.new_page()
            y, margin, lh = 50, 50, 14
            page.insert_text((margin, y), f"Slide {slide_num}", fontsize=16, color=(0, 0, 0.6))
            y += lh * 2

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if y + lh > page.rect.height - margin:
                        page = pdf.new_page(); y = 50
                    page.insert_text((margin, y), shape.text, fontsize=12)
                    y += lh

        pdf.save(output_path)
        pdf.close()
        return output_path
    except Exception as e:
        print(f"[pptx_to_pdf] {e}")
        return None

async def pdf_to_docx(input_path: str, output_path: str) -> str | None:
    if not LIBRE_PATH:
        return None
    try:
        out_dir = str(Path(output_path).parent)
        result = subprocess.run(
            [LIBRE_PATH, "--headless", "--infilter=writer_pdf_import",
             "--convert-to", "docx", "--outdir", out_dir, input_path],
            capture_output=True, timeout=120,
        )
        if result.returncode == 0:
            expected = Path(out_dir) / f"{Path(input_path).stem}.docx"
            if expected.exists():
                expected.rename(output_path)
                return output_path
        return None
    except Exception as e:
        print(f"[pdf_to_docx] {e}")
        return None

async def pdf_to_xlsx(input_path: str, output_path: str) -> str | None:
    if LIBRE_PATH:
        try:
            out_dir = str(Path(output_path).parent)
            result = subprocess.run(
                [LIBRE_PATH, "--headless", "--infilter=writer_pdf_import",
                 "--convert-to", "xlsx", "--outdir", out_dir, input_path],
                capture_output=True, timeout=120,
            )
            if result.returncode == 0:
                expected = Path(out_dir) / f"{Path(input_path).stem}.xlsx"
                if expected.exists():
                    expected.rename(output_path)
                    return output_path
        except Exception:
            pass
    from openpyxl import Workbook
    try:
        doc = fitz.open(input_path)
        wb = Workbook()
        ws = wb.active
        ws.title = "PDF Content"
        row_num = 1
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            for line in text.split("\n"):
                ws.cell(row=row_num, column=1, value=line)
                row_num += 1
            row_num += 1
        wb.save(output_path)
        wb.close()
        doc.close()
        return output_path
    except Exception as e:
        print(f"[pdf_to_xlsx] {e}")
        return None

async def pdf_to_csv(input_path: str, output_path: str) -> str | None:
    import csv
    try:
        doc = fitz.open(input_path)
        with open(output_path, "w", newline="", encoding="utf-8") as f:
            w = csv.writer(f)
            for page_num in range(len(doc)):
                page = doc[page_num]
                text = page.get_text()
                for line in text.split("\n"):
                    w.writerow([line])
        doc.close()
        return output_path
    except Exception as e:
        print(f"[pdf_to_csv] {e}")
        return None
